from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
import io
import os
from typing import List
from models import QuizQuestion

def create_powerpoint(summary: List[str], quiz: List[QuizQuestion]) -> bytes:
    """Create PowerPoint presentation"""
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI Lesson Converter"
    subtitle.text = "Generated Lesson Summary & Quiz"
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]  # Content slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Lesson Summary"
    
    summary_text = ""
    for i, point in enumerate(summary, 1):
        summary_text += f"{i}. {point}\n\n"
    
    content.text = summary_text
    
    # Quiz slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Quiz Questions"
    
    quiz_text = ""
    for i, q in enumerate(quiz, 1):
        quiz_text += f"Q{i}: {q.question}\n"
        for j, option in enumerate(q.options):
            quiz_text += f"  {chr(65+j)}) {option}\n"
        quiz_text += f"Correct: {q.correct_answer}\n\n"
    
    content.text = quiz_text
    
    # Save to bytes
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    return ppt_bytes.getvalue()

def create_pdf(summary: List[str], quiz: List[QuizQuestion]) -> bytes:
    """Create PDF document"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title = Paragraph("AI Lesson Converter", styles['Title'])
    story.append(title)
    story.append(Spacer(1, 0.5*inch))
    
    # Summary section
    summary_title = Paragraph("Lesson Summary", styles['Heading1'])
    story.append(summary_title)
    story.append(Spacer(1, 0.2*inch))
    
    for i, point in enumerate(summary, 1):
        bullet = Paragraph(f"{i}. {point}", styles['Normal'])
        story.append(bullet)
        story.append(Spacer(1, 0.1*inch))
    
    story.append(Spacer(1, 0.3*inch))
    
    # Quiz section
    quiz_title = Paragraph("Quiz Questions", styles['Heading1'])
    story.append(quiz_title)
    story.append(Spacer(1, 0.2*inch))
    
    for i, q in enumerate(quiz, 1):
        question = Paragraph(f"Q{i}: {q.question}", styles['Heading2'])
        story.append(question)
        
        for j, option in enumerate(q.options):
            opt = Paragraph(f"    {chr(65+j)}) {option}", styles['Normal'])
            story.append(opt)
        
        correct = Paragraph(f"    <b>Correct Answer: {q.correct_answer}</b>", styles['Normal'])
        story.append(correct)
        story.append(Spacer(1, 0.2*inch))
    
    doc.build(story)
    buffer.seek(0)
    
    return buffer.getvalue()