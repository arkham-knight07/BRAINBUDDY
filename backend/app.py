from fastapi import FastAPI, HTTPException, Depends, UploadFile, File, status, BackgroundTasks
from fastapi.responses import Response, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import os
from dotenv import load_dotenv
load_dotenv()
import PyPDF2
import logging
import io
from database import init_db, create_user, get_user_by_email, update_user_password
from auth import hash_password, verify_password, create_access_token, verify_token
from models import *
from ai_service import generate_summary, generate_quiz, answer_question
from export_service import create_powerpoint, create_pdf

app = FastAPI(title="BRAINBUDDY API", version="1.0.0", description="API for converting lessons using BRAINBUDDY")

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s %(levelname)s %(message)s')
logger = logging.getLogger(__name__)

# CORS middleware (add production domain)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:5173", "https://your-production-domain.com"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize database on startup
@app.on_event("startup")
def startup_event():
    init_db()

# Store uploaded content temporarily (in production, use proper storage)
uploaded_content = {}

@app.post("/register")
async def register(user_data: UserCreate):
    """Register new user"""
    existing_user = get_user_by_email(user_data.email)
    if existing_user:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Email already registered"
        )
    
    password_hash = hash_password(user_data.password)
    user_id = create_user(user_data.email, password_hash)
    
    return {"message": "User created successfully", "user_id": user_id}

@app.post("/token")
async def login(user_data: UserLogin):
    """Login user and return JWT token"""
    user = get_user_by_email(user_data.email)
    if not user or not verify_password(user_data.password, user["password_hash"]):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Invalid email or password"
        )
    
    access_token = create_access_token({"sub": user_data.email})
    return Token(access_token=access_token)

@app.post("/logout")
async def logout():
    """Logout user by removing token from client (frontend should clear localStorage)"""
    return JSONResponse(content={"message": "Logged out successfully. Please clear your token on the frontend."})



# Forgot password endpoint: generate token and send email (mock)
from auth import generate_reset_token, reset_user_password

def send_reset_email(email: str, token: str):
    # In production, send an actual email
    logger.info(f"Sending password reset email to {email} with token: {token}")

@app.post("/forgot-password")
async def forgot_password(data: dict, background_tasks: BackgroundTasks):
    """Initiate password reset: generate token and send email (mock)"""
    email = data.get("email")
    user = get_user_by_email(email)
    # Always return success for security
    if user:
        token = generate_reset_token(email)
        background_tasks.add_task(send_reset_email, email, token)
    return {"message": "If this email exists, a reset link has been sent."}

# Reset password endpoint
@app.post("/reset-password")
async def reset_password(data: dict):
    """Reset password using token"""
    email = data.get("email")
    token = data.get("token")
    new_password = data.get("new_password")
    if not (email and token and new_password):
        raise HTTPException(status_code=400, detail="Missing fields")
    success = reset_user_password(email, token, new_password)
    if not success:
        raise HTTPException(status_code=400, detail="Invalid token or email")
    return {"message": "Password reset successful"}
@app.post("/upload")
async def upload_file(file: UploadFile = File(...), current_user=Depends(verify_token)):
    try:
        content = ""
        if file.content_type == "application/pdf":
            pdf_bytes = await file.read()
            try:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
            except Exception as e:
                logger.error(f"PDF parsing error: {e}")
                raise HTTPException(status_code=400, detail="Failed to parse PDF file.")
        elif file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            docx_bytes = await file.read()
            try:
                from docx import Document
                doc = Document(io.BytesIO(docx_bytes))
                for para in doc.paragraphs:
                    content += para.text + "\n"
            except Exception as e:
                logger.error(f"DOCX parsing error: {e}")
                raise HTTPException(status_code=400, detail="Failed to parse DOCX file.")
        elif file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            pptx_bytes = await file.read()
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(pptx_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
            except Exception as e:
                logger.error(f"PPTX parsing error: {e}")
                raise HTTPException(status_code=400, detail="Failed to parse PPTX file.")
        elif file.content_type.startswith("text/"):
            text_bytes = await file.read()
            try:
                content = text_bytes.decode("utf-8")
            except Exception as e:
                logger.error(f"Text file decoding error: {e}")
                raise HTTPException(status_code=400, detail="Failed to decode text file.")
        else:
            logger.warning(f"Unsupported file type: {file.content_type}")
            raise HTTPException(status_code=400, detail="Unsupported file type. Please upload PDF, PPTX, DOCX, or text files.")
        uploaded_content[current_user["email"]] = content
        logger.info(f"File uploaded for user: {current_user['email']}")
        return UploadResponse(message="File uploaded successfully", content=content[:500] + "..." if len(content) > 500 else content)
    except HTTPException as e:
        logger.error(f"Upload error: {e.detail}")
        raise
    except Exception as e:
        logger.error(f"Unexpected error: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.post("/summarize")
async def summarize_content(current_user=Depends(verify_token)):
    """Generate summary from uploaded content"""
    user_email = current_user["email"]
    
    if user_email not in uploaded_content:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No content uploaded. Please upload a file first."
        )
    
    content = uploaded_content[user_email]
    summary = generate_summary(content)
    
    return SummaryResponse(summary=summary)

@app.post("/generate_quiz")
async def create_quiz(current_user=Depends(verify_token)):
    """Generate quiz from uploaded content"""
    user_email = current_user["email"]
    
    if user_email not in uploaded_content:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No content uploaded. Please upload a file first."
        )
    
    content = uploaded_content[user_email]
    quiz = generate_quiz(content)
    
    return QuizResponse(quiz=quiz)

@app.post("/ask")
async def ask_question(question_data: dict, current_user=Depends(verify_token)):
    """Answer question about uploaded content"""
    user_email = current_user["email"]
    
    if user_email not in uploaded_content:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No content uploaded. Please upload a file first."
        )
    
    content = uploaded_content[user_email]
    question = question_data.get("question", "")
    
    if not question:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Question is required"
        )
    
    answer = answer_question(content, question)
    
    return AskResponse(question=question, answer=answer)

# Store generated data temporarily
user_data = {}

@app.post("/export_ppt")
async def export_powerpoint(current_user=Depends(verify_token)):
    """Export lesson as PowerPoint"""
    user_email = current_user["email"]
    
    if user_email not in uploaded_content:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No content uploaded. Please upload a file first."
        )
    
    # Generate or retrieve summary and quiz
    content = uploaded_content[user_email]
    summary = generate_summary(content)
    quiz = generate_quiz(content)
    
    # Create PowerPoint
    ppt_bytes = create_powerpoint(summary, quiz)
    
    return Response(
        content=ppt_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=lesson.pptx"}
    )

@app.post("/export_pdf")
async def export_pdf(current_user=Depends(verify_token)):
    """Export lesson as PDF"""
    user_email = current_user["email"]
    
    if user_email not in uploaded_content:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No content uploaded. Please upload a file first."
        )
    
    # Generate or retrieve summary and quiz
    content = uploaded_content[user_email]
    summary = generate_summary(content)
    quiz = generate_quiz(content)
    
    # Create PDF
    pdf_bytes = create_pdf(summary, quiz)
    
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=lesson.pdf"}
    )

@app.post("/generate_flashcards")
async def create_flashcards(current_user=Depends(verify_token)):
    """Generate flashcards from uploaded content"""
    user_email = current_user["email"]
    if user_email not in uploaded_content:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No content uploaded. Please upload a file first."
        )
    content = uploaded_content[user_email]
    from ai_service import generate_flashcards
    flashcards = generate_flashcards(content)
    return {"flashcards": [fc.dict() for fc in flashcards]}

@app.get("/")
async def root():
    return {"message": "AI Lesson Converter API", "status": "running"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)